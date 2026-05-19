# 多 Agent 论文辅助系统 Demo
import asyncio
from typing import List, Dict
import openai
import os
import datetime
from pptx import Presentation
from pptx.util import Pt

openai.api_key = os.getenv("OPENAI_API_KEY")

class Agent:
    def __init__(self, name: str):
        self.name = name
        self.task_queue: List[str] = []

    def assign_task(self, task: str):
        self.task_queue.append(task)
        print(f"[{self.name}] 收到任务: {task}")

    async def run_task(self, task: str) -> str:
        prompt = f"你是专业 AI 助手，任务：{task}"
        try:
            response = openai.ChatCompletion.create(
                model="gpt-4",
                messages=[{"role": "user", "content": prompt}],
                temperature=0.3
            )
            result = response.choices[0].message.content
        except Exception as e:
            result = f"任务执行出错: {e}"

        timestamp = datetime.datetime.now().strftime("%Y-%m-%d %H:%M:%S")
        with open(f"logs/{self.name}_log.txt", "a", encoding="utf-8") as f:
            f.write(f"[{timestamp}] 任务: {task}\n结果: {result}\n\n")
        return result

    async def run(self) -> List[str]:
        results = []
        while self.task_queue:
            task = self.task_queue.pop(0)
            result = await self.run_task(task)
            print(f"[{self.name}] 完成任务: {task}")
            results.append(result)
        return results

class MultiAgentSystem:
    def __init__(self):
        self.agents: Dict[str, Agent] = {}

    def add_agent(self, agent: Agent):
        self.agents[agent.name] = agent

    def assign_task_to_agent(self, agent_name: str, task: str):
        if agent_name in self.agents:
            self.agents[agent_name].assign_task(task)
        else:
            print(f"Agent {agent_name} 不存在")

    async def run_all(self):
        tasks = [agent.run() for agent in self.agents.values()]
        results = await asyncio.gather(*tasks)
        return results

def generate_ppt(title: str, sections: List[str], filename: str):
    prs = Presentation()
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title

    for section in sections:
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = section.split("\n")[0]
        body_shape = slide.shapes.placeholders[1]
        tf = body_shape.text_frame
        tf.text = "\n".join(section.split("\n")[1:])
        for paragraph in tf.paragraphs:
            for run in paragraph.runs:
                run.font.size = Pt(18)

    prs.save(filename)
    print(f"PPT 已生成: {filename}")

async def main():
    system = MultiAgentSystem()
    system.add_agent(Agent("文献整理 Agent"))
    system.add_agent(Agent("PPT 生成 Agent"))
    system.add_agent(Agent("论文校验 Agent"))

    system.assign_task_to_agent("文献整理 Agent", "搜索 AI 自动化论文并整理参考文献，生成 BibTeX")
    system.assign_task_to_agent("文献整理 Agent", "根据整理的文献生成文献摘要列表")
    system.assign_task_to_agent("PPT 生成 Agent", "根据整理的文献摘要生成答辩 PPT 内容")
    system.assign_task_to_agent("论文校验 Agent", "检查论文正文引用与参考文献匹配，生成报告")

    all_results = await system.run_all()

    ppt_sections = all_results[1]  # PPT Agent 结果
    if ppt_sections:
        generate_ppt(
            title="毕业设计答辩 PPT",
            sections=ppt_sections if isinstance(ppt_sections, list) else [ppt_sections],
            filename="outputs/答辩PPT.pptx"
        )

if __name__ == "__main__":
    asyncio.run(main())